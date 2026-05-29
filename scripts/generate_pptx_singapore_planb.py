#!/usr/bin/env python3
"""
Generate a PowerPoint deck for Plan B: Market-entry/Expansion in Singapore.
Usage:
  python3 generate_pptx_singapore_planb.py "Your Company" "May 2026" /path/to/output.pptx
If no output path is provided, defaults to /opt/hermes/output/plan-b-singapore-market-entry.pptx
"""
from pptx import Presentation
from pptx.util import Inches
import sys
import os

TEMPLATE_TITLE_LAYOUT = 0  # Title slide layout index (depends on template)
TEMPLATE_BODY_LAYOUT = 1   # Title and Content layout index

slides_content = [
    {
        "title": "Plan B: Market-entry/Expansion for Singapore",
        "subtitle": None,
        "bullets": [
            "Company: {COMPANY}",
            "Date: {DATE}"
        ],
        "notes": None
    },
    {
        "title": "Agenda",
        "bullets": [
            "Why Singapore",
            "Plan B objectives",
            "Market-entry options (Singapore-focused)",
            "Recommended path and rationale",
            "Execution plan & timeline",
            "Regulatory & compliance considerations",
            "Financials & funding",
            "Risks & mitigations",
            "KPIs & success criteria",
            "Next steps",
        ],
        "notes": None
    },
    {
        "title": "Executive summary (Singapore-specific)",
        "bullets": [
            "Plan B focuses on a low-risk, accelerated entry into Singapore with a clearly defined fallback.",
            "Primary entry path: Direct local entity for control, data locality, and scale potential.",
            "Backup path: Channel/partnership model to de-risk early traction.",
            "Timeline: 12–18 months to validate and scale.",
            "Critical success factors: regulatory clarity, strong local partnerships, access to government grants/incentives."
        ],
        "notes": None
    },
    {
        "title": "Why Singapore is a fit (high level)",
        "bullets": [
            "Strategic gateway to Southeast Asia.",
            "Pro-business environment, strong IP protection, stable legal framework.",
            "Robust digital infrastructure; access to grants and incentives (EnterpriseSingapore, EDB).",
            "English-speaking talent pool and multicultural market."
        ],
        "notes": None
    },
    {
        "title": "Plan B objectives for Singapore",
        "bullets": [
            "Rapid market validation with controlled risk.",
            "Minimized up-front CAPEX while testing demand.",
            "Establish local partnerships to accelerate go-to-market.",
            "Achieve regulatory readiness and scalable operations."
        ],
        "notes": None
    },
    {
        "title": "Market-entry options (Singapore-focused)",
        "bullets": [
            "Option A: Direct local entity (full control, scalable)",
            "Option B: Channel/partner model (distributors, system integrators)",
            "Option C: Strategic alliance or joint venture (shared risk, local credibility)",
            "Option D: Pilots with government or corporates (low-risk experiments)",
            "Note: Recommend a primary path with a clearly defined fallback"
        ],
        "notes": None
    },
    {
        "title": "Recommended path and rationale",
        "bullets": [
            "Primary path: Direct Local Entity -> Why: control, data locality, faster scale, direct customer relationships.",
            "Secondary/fallback: Channel/Partnership -> When to switch: regulatory delays, slower-than-expected pilot uptake, need for rapid traction.",
            "Quick wins: select initial verticals, secure pilot customers, establish local ops."
        ],
        "notes": None
    },
    {
        "title": "Market discovery and customer targeting in Singapore",
        "bullets": [
            "Target segments: insert target sectors/industries",
            "Value proposition tailored to Singapore buyers",
            "Early adopters and pilot candidates",
            "Competitive landscape snapshot (local players vs. substitutes)"
        ],
        "notes": None
    },
    {
        "title": "Go-to-market approach (primary path)",
        "bullets": [
            "Sales motion: direct sales, account-based, inside sales",
            "Partnerships and ecosystem: Engage EDB, industry associations, accelerators",
            "Pricing & packaging: Singapore-ready tiers, local currency, tax considerations",
            "Localized marketing: events, content localized for Singapore market"
        ],
        "notes": None
    },
    {
        "title": "Go-to-market approach (fallback path)",
        "bullets": [
            "Channel/partnership structure and enablement",
            "Partner onboarding, incentives, and enablement",
            "Co-marketing and co-sell motions",
            "Channel risk management and governance"
        ],
        "notes": None
    },
    {
        "title": "Regulatory and compliance considerations (Singapore)",
        "bullets": [
            "Data protection: PDPA basics and data handling",
            "Sector-specific licensing/approvals (if applicable)",
            "Employment laws and local hiring considerations",
            "IP protection and contract enforceability",
            "Tax incentives and grants (overview of relevant programs)"
        ],
        "notes": None
    },
    {
        "title": "Operations and supply chain (Singapore)",
        "bullets": [
            "Entity setup, banking, payroll",
            "Local vendors (legal, tax, HR, accounting)",
            "Data/server considerations",
            "Business continuity planning in Singapore context"
        ],
        "notes": None
    },
    {
        "title": "Financial plan and funding (Singapore)",
        "bullets": [
            "Estimated setup costs and ongoing OPEX in Singapore",
            "Revenue projections and profitability timeline",
            "Government grants/incentives offsets",
            "Financing options (VC, corporate partners, grants)"
        ],
        "notes": None
    },
    {
        "title": "Risk management and mitigations",
        "bullets": [
            "Regulatory/approval delays",
            "Market adoption risk and pivot triggers",
            "Talent acquisition/retention risk",
            "Currency/tax risk considerations",
            "Contingency plans to switch to fallback path"
        ],
        "notes": None
    },
    {
        "title": "KPIs and success metrics",
        "bullets": [
            "Leading indicators: pilot uptake, partner signings, inbound leads",
            "Revenue milestones: 6/12/18 months",
            "Time-to-market: setup, approvals",
            "Compliance/security metrics and audit readiness"
        ],
        "notes": None
    },
    {
        "title": "Implementation timeline (high level)",
        "bullets": [
            "Phase 1: Setup & pilots (0–3 months)",
            "Phase 2: Market-entry execution (4–9 months)",
            "Phase 3: Scale & optimize (10–18 months)",
            "Key milestones and owners"
        ],
        "notes": None
    },
    {
        "title": "Next steps",
        "bullets": [
            "Confirm primary path and fallback",
            "Decide target segments and initial pilots",
            "Assign owners for regulatory, legal, and ops tasks",
            "Schedule a follow-up to finalize budget and approvals"
        ],
        "notes": None
    }
]

OUTPUT_DEFAULT = "/opt/hermes/output/plan-b-singapore-market-entry.pptx"

def main():
    if len(sys.argv) < 2:
        out_path = OUTPUT_DEFAULT
        company = "Your Company"
        date_str = "May 2026"
    else:
        company = sys.argv[1]
        date_str = sys.argv[2] if len(sys.argv) > 2 else "May 2026"
        out_path = sys.argv[3] if len(sys.argv) > 3 else OUTPUT_DEFAULT
    prs = Presentation()
    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Plan B: Market-entry/Expansion for Singapore"
    s1.placeholders[1].text = f"{company} · {date_str}"
    # Slide 2: Agenda
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Agenda"
    b = s2.placeholders[1].text_frame
    b.clear()
    for line in [
        "Why Singapore",
        "Plan B objectives",
        "Market-entry options (Singapore-focused)",
        "Recommended path and rationale",
        "Execution plan & timeline",
        "Regulatory & compliance considerations",
        "Financials & funding",
        "Risks & mitigations",
        "KPIs & success criteria",
        "Next steps",
    ]:
        p = b.add_paragraph(); p.text = line; p.level = 0
    # Subsequent slides from slides_content
    for idx in range(2, len(slides_content)):
        pass
    # We'll load content listically as defined below to maintain order
    for item in slides_content[2:]:
        pass
    # Real implementation: loop over slides_content items to create slides

    # We will construct slides from slides_content data structure (already defined above)
    # Start from after agenda (which is slides_content[0] -> title, slides_content[1] -> agenda)
    for entry in slides_content[2:]:
        title = entry["title"]
        bullets = entry["bullets"]
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.clear()
        for btext in bullets:
            p = tf.add_paragraph(); p.text = btext; p.level = 0
        if entry.get("notes"):
            try:
                s.notes_slide.notes_text_frame.text = entry["notes"]
            except Exception:
                pass
    prs.save(out_path)
    print(out_path)

if __name__ == "__main__":
    main()
